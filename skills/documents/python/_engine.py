"""documents skill — pure document engine (extraction, OCR, overlays, page math, compression) (see handlers.py)."""

from __future__ import annotations

import subprocess
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any

import yaml

from knaif.handler_api import HandlerContext
from knaif.planner import _resolve_path

from ._deps import (
    DocumentsDependencyError,
    _require_pikepdf,
    _require_pillow_image,
    _require_pypdf,
    _require_pypdfium2,
    _require_pytesseract,
    _require_tesseract,
)

TEXT_SUFFIXES = {".txt", ".md"}


IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg"}


OFFICE_SUFFIXES = {".docx", ".pptx", ".xlsx"}


PDF_SUFFIX = ".pdf"


def _replace_fixture_placeholder(value: Any, fixture_path: Path) -> Any:
    if isinstance(value, str):
        return value.replace("{fixture}", str(fixture_path))
    if isinstance(value, list):
        return [_replace_fixture_placeholder(item, fixture_path) for item in value]
    if isinstance(value, dict):
        return {
            key: _replace_fixture_placeholder(item, fixture_path) for key, item in value.items()
        }
    return value


def _input_path(args: dict[str, Any], ctx: HandlerContext) -> Path:
    return _resolve_path(args["input"], ctx.root, ctx.sandbox)


def _input_paths(args: dict[str, Any], ctx: HandlerContext) -> list[Path]:
    return [_resolve_path(path, ctx.root, ctx.sandbox) for path in args["inputs"]]


def _output_path(input_path: Path, args: dict[str, Any], ctx: HandlerContext, suffix: str) -> Path:
    raw = args.get("output")
    if raw:
        return _resolve_path(raw, ctx.root, ctx.sandbox)
    if suffix.startswith("."):
        return input_path.with_suffix(suffix)
    return input_path.with_name(f"{input_path.stem}{suffix}")


def _output_dir(args: dict[str, Any], ctx: HandlerContext, default: Path) -> Path:
    raw = args.get("output_dir")
    if raw:
        return _resolve_path(raw, ctx.root, ctx.sandbox)
    return default


def _profile(skill_dir: Path, quality: str) -> dict[str, Any]:
    path = skill_dir / "profiles" / "compress" / f"{quality}.yaml"
    with path.open(encoding="utf-8") as fh:
        return yaml.safe_load(fh) or {}


def _pdf_page_count(path: Path) -> int:
    pikepdf = _require_pikepdf()
    with pikepdf.Pdf.open(path) as pdf:
        return len(pdf.pages)


def _format_for(path: Path) -> str:
    suffix = path.suffix.lower().lstrip(".")
    return suffix or "unknown"


def _read_text_file(path: Path) -> str:
    return path.read_text(encoding="utf-8")


def _page_records_for_text(text: str) -> list[dict[str, Any]]:
    return [{"page": 1, "text": text}]


def _extract_text_records(path: Path, args: dict[str, Any] | None = None) -> list[dict[str, Any]]:
    args = args or {}
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return _page_records_for_text(_read_text_file(path))
    if suffix in IMAGE_SUFFIXES:
        return _ocr_image_text_records(path, args.get("language", "eng"))
    if suffix == PDF_SUFFIX:
        try:
            from pypdf import PdfReader
        except ImportError as exc:  # pragma: no cover - covered when extra absent
            raise DocumentsDependencyError(
                "Install the documents dep group (uv pip install --group documents) to extract PDF text."
            ) from exc

        reader = PdfReader(str(path))
        if reader.is_encrypted:
            raise ValueError(
                "Cannot extract text from encrypted PDF without a password. Unlock it first."
            )
        selected_pages = _parse_pages(args.get("pages"), len(reader.pages))
        records: list[dict[str, Any]] = []
        for page_number in selected_pages:
            page = reader.pages[page_number - 1]
            records.append({"page": page_number, "text": page.extract_text() or ""})
        if args.get("ocr") or not any(record["text"].strip() for record in records):
            return _ocr_pdf_text_records(path, args, selected_pages)
        return records
    if suffix in OFFICE_SUFFIXES:
        return _extract_office_text(path)
    raise ValueError(f"Unsupported document format: {suffix or path.name}")


def _extract_office_text(path: Path) -> list[dict[str, Any]]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        try:
            import docx
        except ImportError as exc:  # pragma: no cover - covered when extra absent
            raise DocumentsDependencyError(
                "Install the documents dep group (uv pip install --group documents) to extract DOCX text."
            ) from exc
        document = docx.Document(str(path))
        return _page_records_for_text("\n".join(p.text for p in document.paragraphs))
    if suffix == ".pptx":
        try:
            from pptx import Presentation
        except ImportError as exc:  # pragma: no cover - covered when extra absent
            raise DocumentsDependencyError(
                "Install the documents dep group (uv pip install --group documents) to extract PPTX text."
            ) from exc
        prs = Presentation(str(path))
        pages = []
        for index, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
            pages.append({"page": index, "text": "\n".join(texts)})
        return pages
    if suffix == ".xlsx":
        try:
            import openpyxl
        except ImportError as exc:  # pragma: no cover - covered when extra absent
            raise DocumentsDependencyError(
                "Install the documents dep group (uv pip install --group documents) to extract XLSX text."
            ) from exc
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        pages = []
        for index, sheet in enumerate(workbook.worksheets, start=1):
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if cell is None else str(cell) for cell in row]
                if any(cells):
                    rows.append("\t".join(cells).rstrip())
            pages.append({"page": index, "text": "\n".join(rows)})
        return pages
    raise ValueError(f"Unsupported office format: {suffix}")


def _ocr_image_text_records(path: Path, language: str = "eng") -> list[dict[str, Any]]:
    _require_tesseract()
    pytesseract = _require_pytesseract()
    Image = _require_pillow_image()
    with Image.open(path) as image:
        text = pytesseract.image_to_string(image, lang=language)
    return [{"page": 1, "text": text}]


def _render_pdf_page_image(pdf: Any, page_number: int, scale: float = 2.0) -> Any:
    bitmap = pdf[page_number - 1].render(scale=scale)
    return bitmap.to_pil().convert("RGB")


def _ocr_pdf_text_records(
    path: Path,
    args: dict[str, Any],
    selected_pages: list[int] | None = None,
) -> list[dict[str, Any]]:
    _require_tesseract()
    pytesseract = _require_pytesseract()
    pdfium = _require_pypdfium2()
    pdf = pdfium.PdfDocument(str(path))
    try:
        pages = selected_pages or _parse_pages(args.get("pages"), len(pdf))
        language = args.get("language", "eng")
        records = []
        for page_number in pages:
            image = _render_pdf_page_image(pdf, page_number)
            try:
                records.append(
                    {
                        "page": page_number,
                        "text": pytesseract.image_to_string(image, lang=language),
                    }
                )
            finally:
                image.close()
        return records
    finally:
        pdf.close()


def _write_ocr_image_pdf(input_path: Path, output_path: Path, language: str) -> int:
    _require_tesseract()
    pytesseract = _require_pytesseract()
    Image = _require_pillow_image()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with Image.open(input_path) as image:
        pdf_bytes = pytesseract.image_to_pdf_or_hocr(image, extension="pdf", lang=language)
    output_path.write_bytes(pdf_bytes)
    return 1


def _write_ocr_pdf(input_path: Path, output_path: Path, language: str) -> int:
    _require_tesseract()
    pytesseract = _require_pytesseract()
    pypdf = _require_pypdf()
    pdfium = _require_pypdfium2()
    pdf = pdfium.PdfDocument(str(input_path))
    writer = pypdf.PdfWriter()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        for page_number in range(1, len(pdf) + 1):
            image = _render_pdf_page_image(pdf, page_number)
            try:
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                    image,
                    extension="pdf",
                    lang=language,
                )
            finally:
                image.close()
            page_reader = pypdf.PdfReader(BytesIO(pdf_bytes))
            writer.add_page(page_reader.pages[0])
        with output_path.open("wb") as fh:
            writer.write(fh)
        return len(pdf)
    finally:
        pdf.close()


def _resolve_endpoint(token: str, total_pages: int, *, default: int) -> int:
    """Resolve one range endpoint to a 1-based page number.

    Accepts a 1-based integer, the position words ``first``/``last``/``end``,
    or an empty string (open-ended range → *default*). The model emits these
    forms for utterances like "drop the last page" or "pages 3 to the end".
    """
    token = token.strip().lower()
    if token == "":
        return default
    if token == "first":
        return 1
    if token in ("last", "end"):
        return total_pages
    try:
        return int(token)
    except ValueError:
        raise ValueError(f"Unrecognized page reference: {token!r}") from None


def _parse_pages(raw: str | None, total_pages: int, *, strict: bool = True) -> list[int]:
    if raw is None or raw == "":
        return list(range(1, total_pages + 1))

    pages: list[int] = []
    for part in raw.split(","):
        part = part.strip()
        if not part:
            continue
        if part.lower() in ("all", "*"):
            # Models emit "all"/"*" for "every page"; expand to the whole document.
            pages.extend(range(1, total_pages + 1))
            continue
        if "-" in part:
            start_s, end_s = part.split("-", 1)
            start = _resolve_endpoint(start_s, total_pages, default=1)
            end = _resolve_endpoint(end_s, total_pages, default=total_pages)
            if start > end:
                raise ValueError(f"Invalid page range: {part}")
            pages.extend(range(start, end + 1))
        else:
            # `part` is non-empty here (empty parts are skipped above), so the
            # `default` is never consulted for a single token.
            pages.append(_resolve_endpoint(part, total_pages, default=1))

    invalid = [page for page in pages if page < 1 or page > total_pages]
    if invalid and strict:
        raise ValueError(f"Page(s) out of range 1-{total_pages}: {invalid}")
    return pages


def _parse_page_range_specs(raw: str, total_pages: int) -> list[tuple[str, list[int]]]:
    specs: list[tuple[str, list[int]]] = []
    for part in raw.split(","):
        label = part.strip()
        if not label:
            continue
        specs.append((label, _parse_pages(label, total_pages)))
    if not specs:
        raise ValueError("At least one page range is required.")
    return specs


def _snippet(text: str, start: int, end: int, radius: int = 40) -> str:
    left = max(0, start - radius)
    right = min(len(text), end + radius)
    return text[left:right].replace("\n", " ").strip()


def _position_xy(
    position: str, width: float, height: float, margin: float = 36
) -> tuple[float, float]:
    horizontal = "center"
    vertical = "center"
    parts = position.split("-")
    if len(parts) == 2:
        vertical, horizontal = parts
    elif len(parts) == 1:
        vertical = parts[0]

    x = {"left": margin, "center": width / 2, "right": width - margin}.get(horizontal, width / 2)
    y = {"top": height - margin, "center": height / 2, "bottom": margin}.get(vertical, height / 2)
    return x, y


def _make_text_overlay(
    text: str,
    *,
    width: float,
    height: float,
    position: str,
    opacity: float = 1.0,
    font_size: int = 36,
) -> Any:
    try:
        from reportlab.pdfgen import canvas
    except ImportError as exc:  # pragma: no cover - covered when extra absent
        raise DocumentsDependencyError(
            "Install the documents dep group (uv pip install --group documents) for PDF text overlays."
        ) from exc

    packet = BytesIO()
    c = canvas.Canvas(packet, pagesize=(width, height))
    if hasattr(c, "setFillAlpha"):
        c.setFillAlpha(opacity)
    c.setFont("Helvetica", font_size)
    x, y = _position_xy(position, width, height)
    c.drawCentredString(x, y, text)
    c.save()
    packet.seek(0)
    pypdf = _require_pypdf()
    return pypdf.PdfReader(packet).pages[0]


def _make_image_overlay(
    image_path: Path,
    *,
    width: float,
    height: float,
    position: str,
    opacity: float = 1.0,
    target_width: float = 96.0,
) -> Any:
    try:
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas
    except ImportError as exc:  # pragma: no cover - covered when extra absent
        raise DocumentsDependencyError(
            "Install the documents dep group (uv pip install --group documents) for PDF image overlays."
        ) from exc

    Image = _require_pillow_image()
    with Image.open(image_path) as image:
        source = image.convert("RGBA")
        if opacity < 1:
            alpha = source.getchannel("A").point(lambda value: int(value * opacity))
            source.putalpha(alpha)
        image_width, image_height = source.size
    scale = target_width / image_width
    draw_width = target_width
    draw_height = image_height * scale
    center_x, center_y = _position_xy(position, width, height)
    x = max(0, min(width - draw_width, center_x - draw_width / 2))
    y = max(0, min(height - draw_height, center_y - draw_height / 2))

    packet = BytesIO()
    c = canvas.Canvas(packet, pagesize=(width, height))
    with TemporaryDirectory() as tmp:
        overlay_png = Path(tmp) / "overlay.png"
        source.save(overlay_png)
        c.drawImage(
            ImageReader(str(overlay_png)), x, y, width=draw_width, height=draw_height, mask="auto"
        )
        c.save()
    packet.seek(0)
    pypdf = _require_pypdf()
    return pypdf.PdfReader(packet).pages[0]


def _write_text_overlay_pdf(
    *,
    input_path: Path,
    output: Path,
    text_for_page: Any,
    position: str,
    opacity: float = 1.0,
    font_size: int = 36,
) -> None:
    pypdf = _require_pypdf()
    reader = pypdf.PdfReader(str(input_path))
    writer = pypdf.PdfWriter()
    for index, page in enumerate(reader.pages, start=1):
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        overlay_text = text_for_page(index)
        overlay = _make_text_overlay(
            overlay_text,
            width=width,
            height=height,
            position=position,
            opacity=opacity,
            font_size=font_size,
        )
        writer.add_page(page)
        writer.pages[-1].merge_page(overlay)
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("wb") as fh:
        writer.write(fh)


def _write_image_overlay_pdf(
    *,
    input_path: Path,
    output: Path,
    image_path: Path,
    position: str,
    opacity: float = 1.0,
) -> None:
    pypdf = _require_pypdf()
    reader = pypdf.PdfReader(str(input_path))
    writer = pypdf.PdfWriter()
    for page in reader.pages:
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        overlay = _make_image_overlay(
            image_path,
            width=width,
            height=height,
            position=position,
            opacity=opacity,
        )
        writer.add_page(page)
        writer.pages[-1].merge_page(overlay)
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("wb") as fh:
        writer.write(fh)


def _lossless_compress(input_path: Path, output_path: Path) -> None:
    pikepdf = _require_pikepdf()
    with pikepdf.Pdf.open(input_path) as pdf:
        pdf.save(
            output_path,
            compress_streams=True,
            recompress_flate=True,
            object_stream_mode=pikepdf.ObjectStreamMode.generate,
        )


def _ghostscript_compress(input_path: Path, output_path: Path, gs: str, pdfsettings: str) -> None:
    subprocess.run(
        [
            gs,
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.4",
            f"-dPDFSETTINGS={pdfsettings}",
            "-dNOPAUSE",
            "-dQUIET",
            "-dBATCH",
            f"-sOutputFile={output_path}",
            str(input_path),
        ],
        check=True,
        capture_output=True,
        text=True,
    )


def _rasterize_compress(
    input_path: Path, output_path: Path, *, dpi: int, jpeg_quality: int
) -> None:
    try:
        import pypdfium2 as pdfium
    except ImportError as exc:  # pragma: no cover - covered when extra absent
        raise DocumentsDependencyError(
            "Install the documents dep group (uv pip install --group documents) for raster PDF compression."
        ) from exc

    Image = _require_pillow_image()
    pdf = pdfium.PdfDocument(str(input_path))
    pages = []
    scale = dpi / 72
    try:
        for index in range(len(pdf)):
            bitmap = pdf[index].render(scale=scale)
            image = bitmap.to_pil().convert("RGB")
            pages.append(image)
        if not pages:
            raise ValueError("Cannot rasterize an empty PDF.")
        first, rest = pages[0], pages[1:]
        first.save(
            output_path,
            "PDF",
            save_all=True,
            append_images=rest,
            resolution=float(dpi),
            quality=jpeg_quality,
        )
    finally:
        for image in pages:
            if isinstance(image, Image.Image):
                image.close()
        pdf.close()

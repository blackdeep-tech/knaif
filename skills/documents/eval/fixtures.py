"""Deterministic tiny document fixtures for the documents skill."""

from __future__ import annotations

from pathlib import Path

FIXTURE_NAMES = (
    "sample.pdf",
    "sample-protected.pdf",
    "sample-scanned.pdf",
    "sample.docx",
    "sample.pptx",
    "sample.xlsx",
    "sample.png",
    "sample.txt",
    "sample.md",
)

_FIXTURE_COMMAND = "uv run python skills/documents/eval/fixtures.py {output}"
FIXTURES: dict[str, str] = dict.fromkeys(FIXTURE_NAMES, _FIXTURE_COMMAND)

FIXTURE_EXTENSIONS: dict[str, str] = {}


def generate_documents_fixtures(out_dir: Path | str) -> dict[str, Path]:
    """Generate tiny deterministic fixtures and return a manifest.

    Text fixtures are dependency-free. Rich fixtures are generated when the
    optional ``the documents dep group (uv pip install --group documents)`` dependencies are installed.
    """

    root = Path(out_dir)
    root.mkdir(parents=True, exist_ok=True)
    txt = root / "sample.txt"
    md = root / "sample.md"
    txt.write_text("Sample document\nInvoice Alpha\nTotal: 42 EUR\n", encoding="utf-8")
    md.write_text("# Sample\n\nAlpha appears in this markdown fixture.\n", encoding="utf-8")
    manifest = {"txt": txt, "md": md}

    try:
        _generate_pdf_fixtures(root, manifest)
        _generate_office_fixtures(root, manifest)
        _generate_image_fixture(root, manifest)
    except ImportError:
        return manifest

    return manifest


def _generate_pdf_fixtures(root: Path, manifest: dict[str, Path]) -> None:
    from PIL import Image, ImageDraw
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas

    pdf = root / "sample.pdf"
    c = canvas.Canvas(str(pdf), pagesize=letter)
    for text in ("Alpha page one", "Beta page two", "Gamma page three"):
        c.drawString(72, 720, text)
        c.showPage()
    c.save()
    manifest["pdf"] = pdf

    # Password-protected copy (password "secret") so unlock_pdf has a real
    # encrypted input to grade against.
    import pikepdf

    protected = root / "sample-protected.pdf"
    with pikepdf.open(pdf) as src:
        src.save(protected, encryption=pikepdf.Encryption(owner="secret", user="secret"))
    manifest["protected_pdf"] = protected

    scanned_png = root / "sample-scan-source.png"
    image = Image.new("RGB", (640, 360), "white")
    draw = ImageDraw.Draw(image)
    draw.text((40, 150), "Scanned image text", fill="black")
    image.save(scanned_png)

    scanned_pdf = root / "sample-scanned.pdf"
    image.save(scanned_pdf, "PDF", resolution=100.0)
    manifest["scanned_pdf"] = scanned_pdf


def _generate_office_fixtures(root: Path, manifest: dict[str, Path]) -> None:
    import docx
    import openpyxl
    from pptx import Presentation

    docx_path = root / "sample.docx"
    document = docx.Document()
    document.add_paragraph("Docx Alpha paragraph")
    document.add_paragraph("Docx Beta paragraph")
    document.save(docx_path)
    manifest["docx"] = docx_path

    pptx_path = root / "sample.pptx"
    presentation = Presentation()
    for text in ("Slide Alpha", "Slide Beta"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title = slide.shapes.title
        if title is not None:
            title.text = text
    presentation.save(pptx_path)
    manifest["pptx"] = pptx_path

    xlsx_path = root / "sample.xlsx"
    workbook = openpyxl.Workbook()
    first = workbook.active
    first.title = "Alpha"
    first.append(["Name", "Value"])
    first.append(["Sheet Alpha", 42])
    second = workbook.create_sheet("Beta")
    second.append(["Name", "Value"])
    second.append(["Sheet Beta", 84])
    workbook.save(xlsx_path)
    manifest["xlsx"] = xlsx_path


def _generate_image_fixture(root: Path, manifest: dict[str, Path]) -> None:
    from PIL import Image, ImageDraw

    png = root / "sample.png"
    image = Image.new("RGB", (320, 180), "white")
    draw = ImageDraw.Draw(image)
    draw.text((24, 72), "PNG fixture", fill="black")
    image.save(png)
    manifest["png"] = png


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        raise SystemExit("Usage: fixtures.py OUTPUT_FILE")
    generate_documents_fixtures(Path(sys.argv[1]).parent)
